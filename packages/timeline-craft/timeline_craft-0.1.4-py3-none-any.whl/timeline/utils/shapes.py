from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt, Inches


def add_text_box(slide, ppt, left, top, width, height):
    """create a clean text box without any template inheritance"""
    text_box = slide.shapes.add_textbox(
        left=int(ppt.slide_width * left),
        top=int(ppt.slide_height * top),
        width=int(ppt.slide_width * width),
        height=int(ppt.slide_height * height),
    )

    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(7.2)
    text_frame.margin_right = Pt(3.6)
    text_frame.margin_top = Pt(3.6)
    text_frame.margin_bottom = Pt(3.6)

    txBody = text_frame._element
    for child in list(txBody):
        if "lstStyle" in child.tag:
            txBody.remove(child)

    for paragraph in text_frame.paragraphs:
        pPr = paragraph._element.get_or_add_pPr()
        for child in list(pPr):
            pPr.remove(child)

        buNone = OxmlElement("a:buNone")
        pPr.insert(0, buNone)

        pPr.set("marL", "0")
        pPr.set("marR", "0")
        pPr.set("indent", "0")
        pPr.set("lvl", "0")
        pPr.set("algn", "l")

    return text_box


def send_backwards(slide, element):
    slide.shapes._spTree.remove(element._element)
    slide.shapes._spTree.insert(2, element._element)


def SubElement(parent, tagname, **kwargs):
    """create a new element and append it to the parent element"""
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def set_shape_transparency(shape, alpha):
    """set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, "a:alpha", val=str(alpha))
