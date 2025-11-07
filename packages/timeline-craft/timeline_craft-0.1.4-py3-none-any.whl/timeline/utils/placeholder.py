from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
import collections
import collections.abc
from pptx.shapes import shapetree


class _BaseShapes(shapetree._BaseShapes):
    def clone_placeholder(self, placeholder, input_idx=None):
        """add a new placeholder shape based on *placeholder*."""
        sp = placeholder.element
        ph_type, orient, sz, idx = (sp.ph_type, sp.ph_orient, sp.ph_sz, sp.ph_idx)
        if input_idx is not None:
            idx = input_idx
        id_ = self._next_shape_id
        name = self._next_ph_name(ph_type, id_, orient)
        self._spTree.add_placeholder(id_, name, ph_type, orient, sz, idx)
        for shape in self:
            if shape.shape_id == id_:
                return shape


shapetree._BaseShapes.clone_placeholder = _BaseShapes.clone_placeholder


def collect_placeholder_templates(presentation):
    """collect all available placeholder templates in the master layouts"""
    templates = {}
    for layout in presentation.slide_layouts:
        for placeholder in layout.placeholders:
            if placeholder.element.ph_type.__str__() not in templates.keys():
                templates[placeholder.element.ph_type.__str__()] = placeholder

    templates = collections.OrderedDict(sorted(templates.items(), key=lambda x: int("".join(filter(str.isdigit, x[0])))))
    templates = collections.OrderedDict([(key.split(" ")[0], value) for key, value in templates.items()])
    return templates


def add_placeholder(ppt, slide_index: int, template: str, left, top, width, height):
    """add a placeholder to a slide using a template:
        TITLE
        BODY
        CENTER_TITLE
        SUBTITLE
        OBJECT
        SLIDE_NUMBER
        FOOTER
        DATE
        PICTURE
    params:
        slide_index: the index of the slide
        template: the template to use
        left: the left starting coordinate of the placeholder (0-1)
        top: the top starting coordinate of the placeholder (0-1)
        width: the width of the placeholder (0-1)
        height: the height of the placeholder (0-1)
    note:
        if width or height is None, the placeholder will be set as
        symmetric to the left and top coordinates
    """
    slide = ppt.slides[slide_index]

    placeholder_templates = collect_placeholder_templates(ppt)

    num_existing = len(slide.shapes.placeholders)
    placeholder = slide.shapes.clone_placeholder(placeholder_templates[template], input_idx=num_existing)

    placeholder.left = Inches(ppt.slide_width.inches * left)
    placeholder.top = Inches(ppt.slide_height.inches * top)
    placeholder.width = Inches(ppt.slide_width.inches * width)
    placeholder.height = Inches(ppt.slide_height.inches * height)

    text_frame = placeholder.text_frame
    text_frame.margin_left = Pt(7.2)
    text_frame.margin_right = Pt(3.6)
    text_frame.margin_top = Pt(3.6)
    text_frame.margin_bottom = Pt(3.6)

    txBody = text_frame._element
    for child in list(txBody):
        if "lstStyle" in child.tag:
            txBody.remove(child)

    for paragraph in text_frame.paragraphs:
        pPr = paragraph._element.get_or_add_pPr()
        for child in list(pPr):
            tag = child.tag.lower()
            if any(x in tag for x in ["bu", "mar", "indent", "spc", "lvl", "defppr", "tablst"]):
                pPr.remove(child)

        buNone = OxmlElement("a:buNone")
        pPr.insert(0, buNone)

        pPr.set("marL", "0")
        pPr.set("marR", "0")
        pPr.set("indent", "0")
        pPr.set("lvl", "0")
        pPr.set("algn", "l")

    return placeholder
