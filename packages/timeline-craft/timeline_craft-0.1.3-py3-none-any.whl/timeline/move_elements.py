import collections
import collections.abc
from pptx.enum.shapes import MSO_SHAPE
from timeline.utils.shapes import set_shape_transparency, send_backwards, add_text_box
from timeline.utils.paragraph import add_paragraph, amend_font
from timeline.configurations import Configurations


def move_elements_to_right(ppt, config=None):
    """moves all placeholders and images to the right of the screen, shrinking them while maintaining their original aspect ratio, and ensuring their vertical centers remain in the same position"""
    config = config or Configurations()

    for slide in ppt.slides:
        for shape in slide.shapes:
            # original dimensions and positions
            original_top = shape.top
            original_left = shape.left
            original_height = shape.height
            original_width = shape.width

            # calculate scale factor for width
            content_space_width = ppt.slide_width * (1 - config.sidebar_width)
            scale_factor = content_space_width / ppt.slide_width

            # new dimensions
            new_width = original_width * scale_factor
            new_height = original_height * scale_factor

            # new positions
            new_left = (
                ppt.slide_width * config.sidebar_width + (original_left - original_width / 2) * scale_factor + new_width / 2
            )
            vertical_center_offset = (original_height - new_height) / 2
            new_top = original_top + vertical_center_offset

            # apply new dimensions and positions
            shape.left, shape.width = int(new_left), int(new_width)
            shape.top, shape.height = int(new_top), int(new_height)


def merge_tags(tags: list[str]) -> list[str]:
    """merge the adjacent slides with the same tag"""
    merged_tags = []
    for i, tag in enumerate(tags):
        if i == 0:
            merged_tags.append(tag)
        elif tag != merged_tags[-1]:
            merged_tags.append(tag)
    return merged_tags


def set_sidebar_timeline(ppt, tags: list[str], config=None):
    """create sidebar timeline on all slides with the given tags"""
    config = config or Configurations()

    assert len(tags) == len(ppt.slides), f"The number of tags: {len(tags)} has to match the number of slides: {len(ppt.slides)}"

    merged_tags = merge_tags(tags)

    # adding the base shapes to each slide
    for slide in ppt.slides:
        # ------------------------ shaping the sidebar itself ------------------------
        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ppt.slide_width * config.sidebar_width, ppt.slide_height)
        sidebar.name = "!!SIDEBAR"
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = config.sidebar_color
        sidebar.line.color.rgb = config.sidebar_color_outline
        set_shape_transparency(sidebar, config.sidebar_transparency)

        offset = 0
        for tag in merged_tags:
            text_box = add_text_box(
                slide=slide,
                ppt=ppt,
                left=0,
                top=offset,
                width=config.sidebar_width,
                height=config.sidebar_item_height,
            )

            add_paragraph(
                placeholder=text_box,
                text=tag,
                font_size=config.sidebar_init_font_size,
                font_family=config.sidebar_item_font,
                font_color=config.sidebar_item_font_color,
            )

            setattr(text_box, "name", f"!!SIDEBAR_{merged_tags.index(tag)}")

            if tags[ppt.slides.index(slide)] == tag:
                amend_font(
                    placeholder=text_box,
                    font_family=config.sidebar_item_font,
                    font_size=config.sidebar_init_font_size + 3,
                    bold=True,
                )

                indicator = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=0,
                    top=offset * ppt.slide_height,
                    width=ppt.slide_width * (config.sidebar_width + 0.01),
                    height=ppt.slide_height * config.sidebar_item_height,
                )

                indicator.name = "!!INDICATOR"
                indicator.fill.solid()
                indicator.fill.fore_color.rgb = config.indicator_color
                indicator.line.color.rgb = config.sidebar_color_outline
                set_shape_transparency(sidebar, config.indicator_transparency)
                send_backwards(slide, indicator)

            offset += config.sidebar_item_height

        send_backwards(slide, sidebar)


def set_morph_transitions(ppt, config=None):
    """apply morph transitions to all slides with configurable duration"""
    config = config or Configurations()

    for slide in ppt.slides:
        from lxml.etree import Element, SubElement, QName

        sld = slide.element

        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        ns_mc = "http://schemas.openxmlformats.org/markup-compatibility/2006"
        ns_p14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
        ns_p172 = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"

        existing_alt = sld.find(QName(ns_mc, "AlternateContent"))
        if existing_alt is not None:
            sld.remove(existing_alt)

        existing_trans = sld.find(QName(ns_p, "transition"))
        if existing_trans is not None:
            sld.remove(existing_trans)

        cSld = sld.find(QName(ns_p, "cSld"))
        clrMapOvr = sld.find(QName(ns_p, "clrMapOvr"))

        if clrMapOvr is not None:
            idx = list(sld).index(clrMapOvr) + 1
        elif cSld is not None:
            idx = list(sld).index(cSld) + 1
        else:
            continue

        nsmap = {"mc": ns_mc, "p": ns_p, "p14": ns_p14, "p172": ns_p172}

        alt_content = Element(QName(ns_mc, "AlternateContent"), nsmap=nsmap)

        choice = SubElement(alt_content, QName(ns_mc, "Choice"))
        choice.set("Requires", "p159")

        trans_choice = SubElement(choice, QName(ns_p, "transition"))
        trans_choice.set(QName(ns_p14, "dur"), str(int(config.transition_duration * 1000)))

        morph = SubElement(trans_choice, QName(ns_p172, "morph"))
        morph.set("option", "byObject")

        fallback = SubElement(alt_content, QName(ns_mc, "Fallback"))
        trans_fallback = SubElement(fallback, QName(ns_p, "transition"))
        fade = SubElement(trans_fallback, QName(ns_p, "fade"))

        sld.insert(idx, alt_content)
