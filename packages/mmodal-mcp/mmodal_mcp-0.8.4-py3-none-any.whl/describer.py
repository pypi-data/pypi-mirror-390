import asyncio
import base64
import json
import mimetypes
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Tuple

from docx import Document
from litellm import acompletion
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from config import settings


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
PDF_EXTENSIONS = {".pdf"}
DOC_EXTENSIONS = {".docx"}
PPT_EXTENSIONS = {".pptx"}
XLS_EXTENSIONS = {".xlsx"}


@dataclass
class AssetSummary:
    summary: str
    metadata: Dict[str, Any]


def resolve_asset_path(uri: str) -> Path:
    path = Path(uri)
    if path.is_absolute():
        return path.expanduser().resolve()

    image_root = Path(settings.image_dir).expanduser().resolve()

    normalized = path
    try:
        normalized = path.relative_to(image_root.name)
    except ValueError:
        # Path does not start with the image directory name; use as-is.
        pass

    candidate = (image_root / normalized).resolve()
    if candidate.exists():
        return candidate

    # Fallback to original relative path resolution if the above doesn't exist yet.
    return (Path.cwd() / path).resolve()


def _read_image(path: Path) -> Tuple[Dict[str, Any], str]:
    with Image.open(path) as img:
        buffered_metadata = {
            "type": "image",
            "format": img.format,
            "mode": img.mode,
            "width": img.width,
            "height": img.height,
        }

    with path.open("rb") as f:
        encoded = base64.b64encode(f.read()).decode("utf-8")

    mime_type, _ = mimetypes.guess_type(path.name)
    data_url = f"data:{mime_type or 'application/octet-stream'};base64,{encoded}"

    return buffered_metadata, data_url


def _read_pdf(path: Path) -> Tuple[Dict[str, Any], str]:
    reader = PdfReader(str(path))
    text_parts = []
    for page in reader.pages[:5]:
        extracted = page.extract_text() or ""
        text_parts.append(extracted.strip())

    metadata = {
        "type": "pdf",
        "pages": len(reader.pages),
    }

    return metadata, "\n\n".join(part for part in text_parts if part)[:4000]


def _read_docx(path: Path) -> Tuple[Dict[str, Any], str]:
    document = Document(str(path))
    paragraphs = [para.text.strip() for para in document.paragraphs if para.text.strip()]
    metadata = {
        "type": "docx",
        "paragraphs": len(paragraphs),
    }

    return metadata, "\n\n".join(paragraphs)[:4000]


def _read_txt(path: Path) -> Tuple[Dict[str, Any], str]:
    content = path.read_text(encoding="utf-8", errors="ignore")
    metadata = {
        "type": "txt",
        "bytes": path.stat().st_size,
    }

    return metadata, content[:4000]


def _read_pptx(path: Path) -> Tuple[Dict[str, Any], str]:
    presentation = Presentation(str(path))
    slides_text = []
    for index, slide in enumerate(presentation.slides[:10], start=1):
        fragments = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    fragments.append(text)
        if fragments:
            slides_text.append(f"Slide {index}: " + " \n".join(fragments))

    metadata = {
        "type": "pptx",
        "slides": len(presentation.slides),
    }

    return metadata, "\n\n".join(slides_text)[:4000]


def _read_xlsx(path: Path) -> Tuple[Dict[str, Any], str]:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    sheet_summaries = []
    metadata = {
        "type": "xlsx",
        "sheets": workbook.sheetnames,
    }

    for sheet_name in workbook.sheetnames[:5]:
        sheet = workbook[sheet_name]
        rows = []
        row_count = 0
        for row in sheet.iter_rows(min_row=1, max_row=5, max_col=5, values_only=True):
            row_count += 1
            rows.append(
                ", ".join("" if cell is None else str(cell) for cell in row)
            )
        sheet_summaries.append(
            f"Sheet {sheet_name} (sample {min(row_count, 5)} rows):\n" + "\n".join(rows)
        )

    return metadata, "\n\n".join(sheet_summaries)[:4000]


def extract_asset_payload(path: Path) -> Tuple[Dict[str, Any], Dict[str, Any]]:
    ext = path.suffix.lower()

    if ext in IMAGE_EXTENSIONS:
        metadata, data_url = _read_image(path)
        return metadata, {"image_data_url": data_url}
    if ext in PDF_EXTENSIONS:
        metadata, snippet = _read_pdf(path)
        return metadata, {"text_snippet": snippet}
    if ext == ".txt":
        metadata, snippet = _read_txt(path)
        return metadata, {"text_snippet": snippet}
    if ext in DOC_EXTENSIONS:
        metadata, snippet = _read_docx(path)
        return metadata, {"text_snippet": snippet}
    if ext in PPT_EXTENSIONS:
        metadata, snippet = _read_pptx(path)
        return metadata, {"text_snippet": snippet}
    if ext in XLS_EXTENSIONS:
        metadata, snippet = _read_xlsx(path)
        return metadata, {"text_snippet": snippet}

    raise ValueError(f"Unsupported file extension: {ext}")


async def describe_asset(
    uri: str,
    purpose: str = "",
    audience: str = "",
    structure_detail: bool = False,
) -> AssetSummary:
    path = resolve_asset_path(uri)
    if not path.exists():
        raise FileNotFoundError(f"Asset not found at: {path}")

    metadata, payload = await asyncio.to_thread(extract_asset_payload, path)

    user_sections = [settings.asset_description_prompt_template.strip()]
    if purpose:
        user_sections.append(f"Purpose: {purpose.strip()}")
    if audience:
        user_sections.append(f"Audience: {audience.strip()}")
    user_sections.append(f"Metadata: {json.dumps(metadata, indent=2)}")

    content_snippet = payload.get("text_snippet")
    image_data_url = payload.get("image_data_url")
    if content_snippet:
        user_sections.append("Extracted Content Preview:\n" + content_snippet)

    if structure_detail:
        asset_type = metadata.get("type", "")
        if asset_type == "image":
            user_sections.append(settings.asset_structure_guidance_visual.strip())
        else:
            user_sections.append(settings.asset_structure_guidance_document.strip())

    messages: list[Dict[str, Any]] = [
        {"role": "system", "content": settings.asset_description_system_prompt.strip()},
    ]

    if image_data_url:
        user_messages = [
            {
                "type": "text",
                "text": "\n\n".join(user_sections),
            },
            {
                "type": "image_url",
                "image_url": {"url": image_data_url},
            },
        ]
        messages.append({"role": "user", "content": user_messages})
    else:
        messages.append(
            {
                "role": "user",
                "content": "\n\n".join(user_sections),
            }
        )

    llm_settings = settings.get_llm_settings("docs")
    request_kwargs: Dict[str, Any] = {
        "model": llm_settings.model,
        "messages": messages,
    }

    if llm_settings.api_key:
        request_kwargs["api_key"] = llm_settings.api_key
    if llm_settings.api_base:
        request_kwargs["api_base"] = llm_settings.api_base
    if llm_settings.extra_params:
        request_kwargs.update(llm_settings.extra_params)

    response = await acompletion(**request_kwargs)
    choice = response["choices"][0]["message"]["content"]

    return AssetSummary(summary=choice, metadata={**metadata, "path": str(path)})
