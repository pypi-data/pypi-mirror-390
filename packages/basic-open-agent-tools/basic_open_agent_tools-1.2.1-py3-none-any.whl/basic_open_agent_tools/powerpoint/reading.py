"""PowerPoint reading and extraction functions for AI agents."""

import os

from ..decorators import strands_tool

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB limit


@strands_tool
def get_pptx_metadata(file_path: str) -> dict[str, str]:
    """Extract metadata from PowerPoint presentation.


    Args:
        file_path: Path to PowerPoint file (.pptx)

    Returns:
        Dictionary with keys: title, author, subject, keywords, comments, created, modified

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file too large or unreadable

    Example:
        >>> metadata = get_pptx_metadata("/path/to/presentation.pptx")
        >>> metadata['title']
        'My Presentation'
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    file_size = os.path.getsize(file_path)
    if file_size > MAX_FILE_SIZE:
        raise ValueError(
            f"File too large: {file_size} bytes (max {MAX_FILE_SIZE} bytes)"
        )

    try:
        prs = Presentation(file_path)
        core_props = prs.core_properties

        return {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "comments": core_props.comments or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
        }

    except Exception as e:
        raise ValueError(f"Failed to read PowerPoint metadata: {e}")


@strands_tool
def get_pptx_slide_count(file_path: str) -> int:
    """Get number of slides in PowerPoint presentation.

    Args:
        file_path: Path to PowerPoint file (.pptx)

    Returns:
        Number of slides

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file too large or unreadable

    Example:
        >>> count = get_pptx_slide_count("/path/to/presentation.pptx")
        >>> count
        10
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    file_size = os.path.getsize(file_path)
    if file_size > MAX_FILE_SIZE:
        raise ValueError(
            f"File too large: {file_size} bytes (max {MAX_FILE_SIZE} bytes)"
        )

    try:
        prs = Presentation(file_path)
        return len(prs.slides)

    except Exception as e:
        raise ValueError(f"Failed to read PowerPoint slide count: {e}")


@strands_tool
def extract_pptx_text(file_path: str) -> str:
    """Extract all text content from PowerPoint presentation.

    Args:
        file_path: Path to PowerPoint file (.pptx)

    Returns:
        All text content from all slides

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file too large or unreadable

    Example:
        >>> text = extract_pptx_text("/path/to/presentation.pptx")
        >>> len(text) > 0
        True
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    file_size = os.path.getsize(file_path)
    if file_size > MAX_FILE_SIZE:
        raise ValueError(
            f"File too large: {file_size} bytes (max {MAX_FILE_SIZE} bytes)"
        )

    try:
        prs = Presentation(file_path)
        text_parts: list[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)

        return "\n\n".join(text_parts)

    except Exception as e:
        raise ValueError(f"Failed to extract PowerPoint text: {e}")


@strands_tool
def get_pptx_slide_text(file_path: str, slide_index: int) -> str:
    """Get text content from specific slide (0-indexed).

    Args:
        file_path: Path to PowerPoint file (.pptx)
        slide_index: 0-based slide index

    Returns:
        Text content from specified slide

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file too large, unreadable, or slide_index out of range

    Example:
        >>> text = get_pptx_slide_text("/path/to/presentation.pptx", 0)
        >>> len(text) > 0
        True
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")
    if not isinstance(slide_index, int):
        raise TypeError("slide_index must be an integer")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    file_size = os.path.getsize(file_path)
    if file_size > MAX_FILE_SIZE:
        raise ValueError(
            f"File too large: {file_size} bytes (max {MAX_FILE_SIZE} bytes)"
        )

    try:
        prs = Presentation(file_path)

        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(
                f"slide_index {slide_index} out of range (0-{len(prs.slides) - 1})"
            )

        slide = prs.slides[slide_index]
        text_parts: list[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

        return "\n".join(text_parts)

    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Failed to read PowerPoint slide text: {e}")


@strands_tool
def get_pptx_slide_titles(file_path: str) -> list[str]:
    """Get titles of all slides in presentation.

    Args:
        file_path: Path to PowerPoint file (.pptx)

    Returns:
        List of slide titles (empty string if no title)

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file too large or unreadable

    Example:
        >>> titles = get_pptx_slide_titles("/path/to/presentation.pptx")
        >>> titles[0]
        'Introduction'
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    file_size = os.path.getsize(file_path)
    if file_size > MAX_FILE_SIZE:
        raise ValueError(
            f"File too large: {file_size} bytes (max {MAX_FILE_SIZE} bytes)"
        )

    try:
        prs = Presentation(file_path)
        titles: list[str] = []

        for slide in prs.slides:
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            titles.append(title)

        return titles

    except Exception as e:
        raise ValueError(f"Failed to read PowerPoint slide titles: {e}")


@strands_tool
def extract_pptx_notes(file_path: str) -> list[str]:
    """Extract speaker notes from all slides.

    Args:
        file_path: Path to PowerPoint file (.pptx)

    Returns:
        List of notes text for each slide (empty string if no notes)

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file too large or unreadable

    Example:
        >>> notes = extract_pptx_notes("/path/to/presentation.pptx")
        >>> notes[0]
        'Remember to introduce yourself'
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    file_size = os.path.getsize(file_path)
    if file_size > MAX_FILE_SIZE:
        raise ValueError(
            f"File too large: {file_size} bytes (max {MAX_FILE_SIZE} bytes)"
        )

    try:
        prs = Presentation(file_path)
        notes_list: list[str] = []

        for slide in prs.slides:
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text
            notes_list.append(notes_text)

        return notes_list

    except Exception as e:
        raise ValueError(f"Failed to extract PowerPoint notes: {e}")
