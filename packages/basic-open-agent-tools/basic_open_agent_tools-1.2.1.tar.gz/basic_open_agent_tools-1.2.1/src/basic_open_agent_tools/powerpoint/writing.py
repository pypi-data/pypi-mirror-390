"""PowerPoint generation and creation functions for AI agents."""

import os

from ..decorators import strands_tool

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]
    Inches = None  # type: ignore[assignment,misc]
    Pt = None  # type: ignore[assignment,misc]


@strands_tool
def create_simple_pptx(
    file_path: str, title: str, subtitle: str, skip_confirm: bool
) -> str:
    """Create simple PowerPoint presentation with title slide.


    Args:
        file_path: Path for new PowerPoint file (.pptx)
        title: Title text for title slide
        subtitle: Subtitle text for title slide
        skip_confirm: If False, raises error if file exists

    Returns:
        Success message

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        ValueError: If file exists and skip_confirm is False

    Example:
        >>> msg = create_simple_pptx("/tmp/pres.pptx", "Title", "Subtitle", True)
        >>> "Created" in msg
        True
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")
    if not isinstance(title, str):
        raise TypeError("title must be a string")
    if not isinstance(subtitle, str):
        raise TypeError("subtitle must be a string")
    if not isinstance(skip_confirm, bool):
        raise TypeError("skip_confirm must be a boolean")

    if os.path.exists(file_path) and not skip_confirm:
        raise ValueError(
            f"File already exists: {file_path}. Set skip_confirm=True to overwrite."
        )

    parent_dir = os.path.dirname(file_path)
    if parent_dir and not os.path.exists(parent_dir):
        raise ValueError(f"Parent directory does not exist: {parent_dir}")

    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

        prs.save(file_path)

        return f"Created PowerPoint presentation at {file_path}"

    except PermissionError:
        raise PermissionError(f"Permission denied writing to: {file_path}")
    except Exception as e:
        raise ValueError(f"Failed to create PowerPoint presentation: {e}")


@strands_tool
def add_pptx_title_slide(file_path: str, title: str, subtitle: str) -> str:
    """Add title slide to existing PowerPoint presentation.

    Args:
        file_path: Path to existing PowerPoint file (.pptx)
        title: Title text for slide
        subtitle: Subtitle text for slide

    Returns:
        Success message

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file unreadable

    Example:
        >>> msg = add_pptx_title_slide("/tmp/pres.pptx", "New Title", "Subtitle")
        >>> "Added" in msg
        True
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")
    if not isinstance(title, str):
        raise TypeError("title must be a string")
    if not isinstance(subtitle, str):
        raise TypeError("subtitle must be a string")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    try:
        prs = Presentation(file_path)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

        prs.save(file_path)

        return f"Added title slide to {file_path}"

    except Exception as e:
        raise ValueError(f"Failed to add title slide: {e}")


@strands_tool
def add_pptx_content_slide(file_path: str, title: str, bullet_points: list[str]) -> str:
    """Add content slide with title and bullet points.

    Args:
        file_path: Path to existing PowerPoint file (.pptx)
        title: Title text for slide
        bullet_points: List of bullet point strings

    Returns:
        Success message

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file unreadable or bullet_points empty

    Example:
        >>> msg = add_pptx_content_slide("/tmp/pres.pptx", "Agenda", ["Point 1", "Point 2"])
        >>> "Added" in msg
        True
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")
    if not isinstance(title, str):
        raise TypeError("title must be a string")
    if not isinstance(bullet_points, list):
        raise TypeError("bullet_points must be a list")
    if not bullet_points:
        raise ValueError("bullet_points must not be empty")

    for point in bullet_points:
        if not isinstance(point, str):
            raise TypeError("all bullet_points must be strings")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    try:
        prs = Presentation(file_path)
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)

        slide.shapes.title.text = title

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.level = 0

        prs.save(file_path)

        return f"Added content slide to {file_path}"

    except Exception as e:
        raise ValueError(f"Failed to add content slide: {e}")


@strands_tool
def add_pptx_blank_slide(file_path: str) -> str:
    """Add blank slide to existing PowerPoint presentation.

    Args:
        file_path: Path to existing PowerPoint file (.pptx)

    Returns:
        Success message

    Raises:
        ImportError: If python-pptx is not installed
        TypeError: If parameters are wrong type
        FileNotFoundError: If file doesn't exist
        ValueError: If file unreadable

    Example:
        >>> msg = add_pptx_blank_slide("/tmp/pres.pptx")
        >>> "Added" in msg
        True
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for PowerPoint support. "
            "Install with: pip install python-pptx"
        )

    if not isinstance(file_path, str):
        raise TypeError("file_path must be a string")

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    try:
        prs = Presentation(file_path)
        blank_slide_layout = prs.slide_layouts[6]
        prs.slides.add_slide(blank_slide_layout)

        prs.save(file_path)

        return f"Added blank slide to {file_path}"

    except Exception as e:
        raise ValueError(f"Failed to add blank slide: {e}")
