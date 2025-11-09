"""
OECT器件稳定性分析报告生成Pipeline v2.0
============================================

重构版本，采用两阶段架构：
1. 素材生成阶段：为每个设备单独生成图片和视频
2. PPT排布阶段：从本地素材读取并排布成专业PPT

改进特性：
- 每个设备的图片单独保存，不再合成大图
- 视频正确嵌入PPT而非仅缩略图
- 删除不必要的最后两页
- 优化的素材管理和PPT布局

使用方法：
    python stability_report_pipeline_v2.py --chip "#20250804008" --output "stability_report.pptx"
    python stability_report_pipeline_v2.py --all-chips --output-dir "reports/"
    python stability_report_pipeline_v2.py --all-chips --output-dir "/home/lidonghaowsl/develop_win/hdd/data/Stability_PS20250929/reports"

作者: OECT Data Processing Team v2.0
"""

import os
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, Union
import argparse
import logging
from datetime import datetime
import tempfile
import json
import warnings
warnings.filterwarnings('ignore')

import numpy as np
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
from matplotlib.figure import Figure
import pandas as pd

# PPT生成相关
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. PPT generation will be disabled.")
    print("Install with: pip install python-pptx")

# 项目模块
try:
    from catalog import UnifiedExperimentManager, find_experiments
    from visualization import OECTPlotter, ChipFeaturePlotter, plot_chip_feature
    from experiment import Experiment
    from features import FeatureReader
except ImportError as e:
    print(f"Error importing project modules: {e}")
    print("Make sure all required modules are available")
    sys.exit(1)

# 配置日志
from ..logger_config import get_module_logger
logger = get_module_logger()


class AssetManager:
    """素材管理器 - 负责组织和管理生成的图片、视频素材"""
    
    def __init__(self, base_dir: Union[str, Path]):
        """
        初始化素材管理器
        
        Args:
            base_dir: 素材存储基础目录
        """
        self.base_dir = Path(base_dir)
        self.base_dir.mkdir(parents=True, exist_ok=True)
        
        # 创建子目录结构
        self.images_dir = self.base_dir / "images"
        self.videos_dir = self.base_dir / "videos"
        self.features_dir = self.base_dir / "features"
        
        for dir_path in [self.images_dir, self.videos_dir, self.features_dir]:
            dir_path.mkdir(parents=True, exist_ok=True)
            
        logger.info(f"Asset manager initialized with base directory: {self.base_dir}")
    
    def get_image_path(self, chip_id: str, device_id: str, image_type: str) -> Path:
        """获取图片文件路径"""
        filename = f"{chip_id}_{device_id}_{image_type}.png"
        return self.images_dir / filename
    
    def get_video_path(self, chip_id: str, device_id: str) -> Path:
        """获取视频文件路径"""
        filename = f"{chip_id}_{device_id}_transfer_evolution.mp4"
        return self.videos_dir / filename
    
    def get_feature_path(self, chip_id: str, feature_name: str, variant: str) -> Path:
        """获取特征图片文件路径"""
        filename = f"{chip_id}_{feature_name}_{variant}.png"
        return self.features_dir / filename
    
    def list_chip_assets(self, chip_id: str) -> Dict[str, List[Path]]:
        """列出指定chip的所有素材文件"""
        assets = {
            'images': list(self.images_dir.glob(f"{chip_id}_*.png")),
            'videos': list(self.videos_dir.glob(f"{chip_id}_*.mp4")),
            'features': list(self.features_dir.glob(f"{chip_id}_*.png"))
        }
        return assets


class AssetGenerator:
    """素材生成器 - 第一阶段：生成所有图片和视频素材"""
    
    def __init__(self, 
                 asset_manager: AssetManager,
                 config_path: str = 'catalog_config.yaml',
                 overwrite_mode: bool = False):
        """
        初始化素材生成器
        
        Args:
            asset_manager: 素材管理器实例
            config_path: catalog配置文件路径
            overwrite_mode: 是否覆盖已存在的文件。True=覆盖，False=跳过已存在的文件（默认）
        """
        self.asset_manager = asset_manager
        self.overwrite_mode = overwrite_mode
        
        # 设备ID列表（固定顺序）
        self.device_ids = ['1', '2', '3', '4', '5', '6']
        
        # 初始化catalog管理器
        try:
            self.manager = UnifiedExperimentManager(config_path)
            logger.info("Catalog manager initialized successfully")
        except Exception as e:
            logger.error(f"Failed to initialize catalog manager: {e}")
            raise
        
        # 特征绘图器（延迟初始化）
        self.feature_plotter = None
    
    def generate_chip_assets(self, chip_id: str) -> Dict[str, Any]:
        """
        为指定chip生成所有素材
        
        Args:
            chip_id: chip ID
            
        Returns:
            Dict[str, Any]: 生成结果统计
        """
        logger.info(f"Starting asset generation for chip {chip_id}")
        
        # 获取chip的所有设备
        devices = self._get_chip_devices(chip_id)
        if not devices:
            logger.error(f"No devices found for chip {chip_id}")
            return {'success': False, 'error': 'No devices found'}
        
        results = {
            'success': True,
            'chip_id': chip_id,
            'devices_found': len(devices),
            'generated_assets': {
                'images': 0,
                'videos': 0,
                'features': 0
            },
            'failed_assets': []
        }
        
        # 1. 生成每个设备的图片
        logger.info(f"Generating individual device images for chip {chip_id}")
        for device_id in self.device_ids:
            if device_id in devices:
                try:
                    # 生成各种类型的图片
                    self._generate_device_images(chip_id, device_id, devices[device_id])
                    results['generated_assets']['images'] += 5  # 5种图片类型
                except Exception as e:
                    logger.error(f"Failed to generate images for device {device_id}: {e}")
                    results['failed_assets'].append(f"images_{device_id}")
        
        # 2. 生成每个设备的视频
        logger.info(f"Generating device videos for chip {chip_id}")
        for device_id in self.device_ids:
            if device_id in devices:
                try:
                    self._generate_device_video(chip_id, device_id, devices[device_id])
                    results['generated_assets']['videos'] += 1
                except Exception as e:
                    logger.error(f"Failed to generate video for device {device_id}: {e}")
                    results['failed_assets'].append(f"video_{device_id}")
        
        # 3. 生成特征分析图片
        logger.info(f"Generating feature analysis images for chip {chip_id}")
        try:
            feature_count = self._generate_feature_images(chip_id)
            results['generated_assets']['features'] = feature_count
        except Exception as e:
            logger.error(f"Failed to generate feature images: {e}")
            results['failed_assets'].append("features")
        
        logger.info(f"Asset generation completed for chip {chip_id}. "
                   f"Generated: {results['generated_assets']}, "
                   f"Failed: {len(results['failed_assets'])}")
        
        return results
    
    def _get_chip_devices(self, chip_id: str) -> Dict[str, Any]:
        """获取指定chip的所有设备信息"""
        logger.info(f"Getting devices for chip {chip_id}...")
        
        try:
            # 搜索该chip的所有实验
            experiments = self.manager.search(chip_id=chip_id)
            
            devices = {}
            for exp in experiments:
                device_id = exp.device_id
                if device_id in self.device_ids:
                    devices[device_id] = {
                        'experiment': exp,
                        'chip_id': exp.chip_id,
                        'device_id': device_id,
                        'test_id': exp.test_id,
                        'status': exp.status,
                        'completion_percentage': exp.completion_percentage
                    }
            
            logger.info(f"Found {len(devices)} devices for chip {chip_id}: {list(devices.keys())}")
            return devices
            
        except Exception as e:
            logger.error(f"Failed to get devices for chip {chip_id}: {e}")
            return {}
    
    def _generate_device_images(self, chip_id: str, device_id: str, device_info: Dict[str, Any]):
        """为单个设备生成所有类型的图片"""
        exp = device_info['experiment']
        exp_file = exp.file_path
        
        if not exp_file or not Path(exp_file).exists():
            logger.warning(f"Experiment file not found for device {device_id}: {exp_file}")
            return
        
        plotter = OECTPlotter(exp_file)
        
        # 图片类型配置
        image_configs = [
            {
                'type': 'transient_all',
                'method': 'plot_transient_all',
                'params': {'figsize': (10, 6)}
            },
            {
                'type': 'transfer_evolution_linear',
                'method': 'plot_transfer_evolution',
                'params': {'figsize': (10, 6), 'log_scale': False}
            },
            {
                'type': 'transfer_evolution_log',
                'method': 'plot_transfer_evolution', 
                'params': {'figsize': (10, 6), 'log_scale': True}
            },
            {
                'type': 'transient_early',
                'method': 'plot_transient_single',
                'params': {'step_index': 1, 'time_range': (5, 6), 'figsize': (10, 6)}
            },
            {
                'type': 'transient_late',
                'method': 'plot_transient_single',
                'params': {'step_index': 999, 'time_range': (5, 6), 'figsize': (10, 6)}
            }
        ]
        
        for config in image_configs:
            try:
                # 获取图片路径并检查是否已存在
                image_path = self.asset_manager.get_image_path(chip_id, device_id, config['type'])
                
                # 如果不是覆盖模式且文件已存在，则跳过
                if not self.overwrite_mode and image_path.exists():
                    logger.info(f"Skipping {config['type']} for device {device_id}: file already exists at {image_path}")
                    continue
                
                # 调用绘图方法
                method = getattr(plotter, config['method'])
                fig = method(**config['params'])
                
                # 保存图片（不修改原有标题）
                if fig:
                    fig.savefig(image_path, dpi=300, bbox_inches='tight')
                    plt.close(fig)
                    
                    logger.info(f"Generated {config['type']} image for device {device_id}: {image_path}")
                
            except Exception as e:
                logger.error(f"Failed to generate {config['type']} for device {device_id}: {e}")
    
    def _generate_device_video(self, chip_id: str, device_id: str, device_info: Dict[str, Any]):
        """为单个设备生成视频"""
        exp = device_info['experiment']
        exp_file = exp.file_path
        
        if not exp_file or not Path(exp_file).exists():
            logger.warning(f"Experiment file not found for device {device_id}: {exp_file}")
            return
        
        try:
            plotter = OECTPlotter(exp_file)
            video_path = self.asset_manager.get_video_path(chip_id, device_id)
            
            # 如果不是覆盖模式且文件已存在，则跳过
            if not self.overwrite_mode and video_path.exists():
                logger.info(f"Skipping video for device {device_id}: file already exists at {video_path}")
                return
            
            # 生成视频
            final_path = plotter.create_transfer_video_parallel(
                save_path=str(video_path),
                fps=30,
                figsize=(10, 6),
                verbose=False  # 减少输出
            )
            
            if final_path and Path(final_path).exists():
                logger.info(f"Generated video for device {device_id}: {final_path}")
            else:
                logger.warning(f"Video generation failed for device {device_id}")
                
        except Exception as e:
            logger.error(f"Failed to generate video for device {device_id}: {e}")
    
    def _init_feature_plotter(self):
        """延迟初始化特征绘图器"""
        if self.feature_plotter is None:
            # 获取features目录路径
            features_dir = "/mnt/d/UserData/lidonghaowsl/data/Stability_PS/features/"
            if not Path(features_dir).exists():
                features_dir = "data/features/"
            
            try:
                self.feature_plotter = ChipFeaturePlotter(features_dir)
                logger.info(f"Feature plotter initialized with directory: {features_dir}")
            except Exception as e:
                logger.warning(f"Failed to initialize feature plotter: {e}")
                self.feature_plotter = None
    
    def _generate_feature_images(self, chip_id: str) -> int:
        """生成特征分析图片"""
        self._init_feature_plotter()
        
        if not self.feature_plotter:
            logger.warning("Feature plotter not available, skipping feature image generation")
            return 0
        
        feature_configs = [
            # absI_max_raw特征的6个变体
            {'feature': 'absI_max_raw', 'skip_points': 0, 'normalize_to_first': False, 'variant': 'raw_data'},
            {'feature': 'absI_max_raw', 'skip_points': 1, 'normalize_to_first': False, 'variant': 'skip_1_point'},
            {'feature': 'absI_max_raw', 'skip_points': 5, 'normalize_to_first': False, 'variant': 'skip_5_points'},
            {'feature': 'absI_max_raw', 'skip_points': 0, 'normalize_to_first': True, 'variant': 'normalized_to_first'},
            {'feature': 'absI_max_raw', 'skip_points': 1, 'normalize_to_first': True, 'variant': 'skip_1_normalized'},
            {'feature': 'absI_max_raw', 'skip_points': 5, 'normalize_to_first': True, 'variant': 'skip_5_normalized'},
            
            # absgm_max_forward特征的6个变体
            {'feature': 'absgm_max_forward', 'skip_points': 0, 'normalize_to_first': False, 'variant': 'raw_data'},
            {'feature': 'absgm_max_forward', 'skip_points': 1, 'normalize_to_first': False, 'variant': 'skip_1_point'},
            {'feature': 'absgm_max_forward', 'skip_points': 5, 'normalize_to_first': False, 'variant': 'skip_5_points'},
            {'feature': 'absgm_max_forward', 'skip_points': 0, 'normalize_to_first': True, 'variant': 'normalized_to_first'},
            {'feature': 'absgm_max_forward', 'skip_points': 1, 'normalize_to_first': True, 'variant': 'skip_1_normalized'},
            {'feature': 'absgm_max_forward', 'skip_points': 5, 'normalize_to_first': True, 'variant': 'skip_5_normalized'}
        ]
        
        generated_count = 0
        
        for config in feature_configs:
            try:
                # 获取特征图片路径并检查是否已存在
                feature_path = self.asset_manager.get_feature_path(
                    chip_id, config['feature'], config['variant']
                )
                
                # 如果不是覆盖模式且文件已存在，则跳过
                if not self.overwrite_mode and feature_path.exists():
                    logger.info(f"Skipping feature image {config['feature']}_{config['variant']}: file already exists at {feature_path}")
                    generated_count += 1  # 计数已存在的文件
                    continue
                
                # 生成特征图
                fig = self.feature_plotter.plot_chip_feature(
                    chip_id=chip_id,
                    feature_name=config['feature'],
                    skip_points=config['skip_points'],
                    normalize_to_first=config['normalize_to_first'],
                    figsize=(12, 8)
                )
                
                if fig:
                    # 保存特征图片
                    fig.savefig(feature_path, dpi=300, bbox_inches='tight')
                    plt.close(fig)
                    
                    logger.info(f"Generated feature image: {feature_path}")
                    generated_count += 1
                
            except Exception as e:
                logger.error(f"Failed to generate feature image for {config['feature']}_{config['variant']}: {e}")
        
        return generated_count


class PPTComposer:
    """PPT编排器 - 第二阶段：从本地素材组装专业PPT"""
    
    def __init__(self, asset_manager: AssetManager):
        """
        初始化PPT编排器
        
        Args:
            asset_manager: 素材管理器实例
        """
        self.asset_manager = asset_manager
        
        # PPT尺寸设置（16:9）
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        
        # 设备ID列表
        self.device_ids = ['1', '2', '3', '4', '5', '6']
        
        # 设备信息（简化版，用于显示）
        self.device_info = {
            '1': {'area': '封装面积最小', 'semiconductor_area': '工作区域最大'},
            '2': {'area': '封装面积较小', 'semiconductor_area': '工作区域较大'},
            '3': {'area': '封装面积中等', 'semiconductor_area': '工作区域中等'},
            '4': {'area': '封装面积较大', 'semiconductor_area': '工作区域较小'},
            '5': {'area': '封装面积很大', 'semiconductor_area': '工作区域很小'},
            '6': {'area': '封装面积最大', 'semiconductor_area': '工作区域最小'}
        }
    
    def create_presentation(self, chip_id: str, output_path: Union[str, Path]) -> str:
        """
        创建完整的PPT演示文稿
        
        Args:
            chip_id: chip ID
            output_path: 输出PPT文件路径
            
        Returns:
            str: 最终PPT文件路径
        """
        if not HAS_PPTX:
            logger.error("python-pptx not available. Cannot create PPT.")
            return None
        
        logger.info(f"Creating PPT presentation for chip {chip_id}")
        
        # 检查素材是否存在
        assets = self.asset_manager.list_chip_assets(chip_id)
        if not any(assets.values()):
            logger.error(f"No assets found for chip {chip_id}")
            return None
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        # 添加标题页
        self._add_title_slide(prs, chip_id)
        
        # 添加内容页面
        self._add_device_grid_slide(prs, "Transient Response - All Devices", 
                                  chip_id, "transient_all", 1)
        
        self._add_device_grid_slide(prs, "Transfer Evolution - Linear Scale", 
                                  chip_id, "transfer_evolution_linear", 2)
        
        self._add_device_grid_slide(prs, "Transfer Evolution - Log Scale", 
                                  chip_id, "transfer_evolution_log", 3)
        
        # 视频页面
        self._add_video_grid_slide(prs, "Transfer Evolution Videos", chip_id, 4)
        
        self._add_device_grid_slide(prs, "Early Transient Response (Step 1)", 
                                  chip_id, "transient_early", 5)
        
        self._add_device_grid_slide(prs, "Late Transient Response (Step 1000)", 
                                  chip_id, "transient_late", 6)
        
        # 特征分析页面
        self._add_feature_grid_slide(prs, "Feature Analysis - absI_max_raw", 
                                   chip_id, "absI_max_raw", 8)
        
        self._add_feature_grid_slide(prs, "Feature Analysis - absgm_max_forward", 
                                   chip_id, "absgm_max_forward", 9)
        
        # 保存PPT
        final_output = Path(output_path)
        final_output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(final_output))
        
        logger.info(f"PPT saved: {final_output}")
        return str(final_output)
    
    def _add_title_slide(self, prs: Presentation, chip_id: str):
        """添加标题幻灯片"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"OECT器件稳定性分析报告"
        subtitle.text = f"Chip {chip_id} - 封装面积对稳定性影响研究\n\n生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    def _add_device_grid_slide(self, prs: Presentation, title: str, chip_id: str, image_type: str, page_num: int):
        """添加设备网格布局幻灯片（2行3列）"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = f"{title} - {chip_id}"
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 计算2行3列布局
        margin = Inches(0.2)
        start_x = margin
        start_y = Inches(1.2)
        available_width = Inches(12.33) - 2 * margin
        available_height = Inches(5.8)
        
        image_width = available_width / 3
        image_height = available_height / 2
        
        # 添加设备图片
        for i, device_id in enumerate(self.device_ids):
            row = i // 3
            col = i % 3
            x = start_x + col * image_width
            y = start_y + row * image_height
            
            # 获取图片路径
            image_path = self.asset_manager.get_image_path(chip_id, device_id, image_type)
            
            if image_path.exists():
                try:
                    # 添加图片
                    slide.shapes.add_picture(str(image_path), x + Inches(0.05), y + Inches(0.05),
                                           width=image_width - Inches(0.1), 
                                           height=image_height - Inches(0.3))
                    
                    # 添加设备标签
                    label_shape = slide.shapes.add_textbox(x, y + image_height - Inches(0.25),
                                                         image_width, Inches(0.2))
                    label_frame = label_shape.text_frame
                    label_frame.text = f"Device {device_id}"
                    label_frame.paragraphs[0].font.size = Pt(12)
                    label_frame.paragraphs[0].font.bold = True
                    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                except Exception as e:
                    logger.warning(f"Failed to add image for device {device_id}: {e}")
                    self._add_placeholder_box(slide, x, y, image_width, image_height, 
                                            f"Device {device_id}\nImage Error")
            else:
                # 添加占位符
                self._add_placeholder_box(slide, x, y, image_width, image_height, 
                                        f"Device {device_id}\nNo Data")
    
    def _add_video_grid_slide(self, prs: Presentation, title: str, chip_id: str, page_num: int):
        """添加视频网格布局幻灯片（直接嵌入视频文件）"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = f"{title} - {chip_id}"
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 布局参数
        margin = Inches(0.2)
        start_x = margin
        start_y = Inches(1.2)
        available_width = Inches(12.33) - 2 * margin
        available_height = Inches(5.8)
        
        video_width = available_width / 3
        video_height = available_height / 2
        
        # 添加设备视频
        for i, device_id in enumerate(self.device_ids):
            row = i // 3
            col = i % 3
            x = start_x + col * video_width
            y = start_y + row * video_height
            
            # 获取视频路径
            video_path = self.asset_manager.get_video_path(chip_id, device_id)
            
            if video_path.exists():
                try:
                    # 直接插入视频文件到PPT
                    left = x + Inches(0.05)
                    top = y + Inches(0.05)
                    width = video_width - Inches(0.1)
                    height = video_height - Inches(0.3)
                    
                    # 添加视频
                    video_shape = slide.shapes.add_movie(str(video_path), left, top, width, height)
                    
                    # 添加视频标签
                    label_shape = slide.shapes.add_textbox(x, y + video_height - Inches(0.25),
                                                         video_width, Inches(0.2))
                    label_frame = label_shape.text_frame
                    label_frame.text = f"Device {device_id} - Transfer Evolution"
                    label_frame.paragraphs[0].font.size = Pt(10)
                    label_frame.paragraphs[0].font.bold = True
                    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    logger.info(f"Added video for device {device_id}: {video_path}")
                    
                except Exception as e:
                    logger.warning(f"Failed to add video for device {device_id}: {e}")
                    self._add_placeholder_box(slide, x, y, video_width, video_height,
                                            f"Device {device_id}\nVideo Error")
            else:
                # 添加占位符
                self._add_placeholder_box(slide, x, y, video_width, video_height,
                                        f"Device {device_id}\nNo Video")
    
    def _add_feature_grid_slide(self, prs: Presentation, title: str, chip_id: str, 
                              feature_name: str, page_num: int):
        """添加特征分析网格布局幻灯片（2行3列，6个变体）"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = f"{title} - {chip_id}"
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 布局参数
        margin = Inches(0.1)
        start_x = margin
        start_y = Inches(1.2)
        available_width = Inches(12.33) - 2 * margin
        available_height = Inches(5.8)
        
        image_width = available_width / 3
        image_height = available_height / 2
        
        # 特征变体配置
        variants = [
            {'variant': 'raw_data', 'title': 'Raw Data'},
            {'variant': 'skip_1_point', 'title': 'Skip 1 Point'},
            {'variant': 'skip_5_points', 'title': 'Skip 5 Points'},
            {'variant': 'normalized_to_first', 'title': 'Normalized to First'},
            {'variant': 'skip_1_normalized', 'title': 'Skip 1 + Normalized'},
            {'variant': 'skip_5_normalized', 'title': 'Skip 5 + Normalized'}
        ]
        
        for i, variant_config in enumerate(variants):
            row = i // 3
            col = i % 3
            x = start_x + col * image_width
            y = start_y + row * image_height
            
            # 获取特征图片路径
            feature_path = self.asset_manager.get_feature_path(
                chip_id, feature_name, variant_config['variant']
            )
            
            if feature_path.exists():
                try:
                    # 添加特征图片
                    slide.shapes.add_picture(str(feature_path), x + Inches(0.05), y + Inches(0.05),
                                           width=image_width - Inches(0.1),
                                           height=image_height - Inches(0.3))
                    
                    # 添加变体标签
                    label_shape = slide.shapes.add_textbox(x, y + image_height - Inches(0.25),
                                                         image_width, Inches(0.2))
                    label_frame = label_shape.text_frame
                    label_frame.text = variant_config['title']
                    label_frame.paragraphs[0].font.size = Pt(10)
                    label_frame.paragraphs[0].font.bold = True
                    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                except Exception as e:
                    logger.warning(f"Failed to add feature image {variant_config['variant']}: {e}")
                    self._add_placeholder_box(slide, x, y, image_width, image_height,
                                            f"{variant_config['title']}\nImage Error")
            else:
                # 添加占位符
                self._add_placeholder_box(slide, x, y, image_width, image_height,
                                        f"{variant_config['title']}\nNo Data")
    
    def _add_placeholder_box(self, slide, x: Inches, y: Inches, 
                           width: Inches, height: Inches, text: str):
        """添加占位符框"""
        placeholder_shape = slide.shapes.add_textbox(x, y, width, height)
        placeholder_frame = placeholder_shape.text_frame
        placeholder_frame.text = text
        placeholder_frame.paragraphs[0].font.size = Pt(12)
        placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 设置样式
        placeholder_shape.fill.solid()
        placeholder_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
        placeholder_shape.line.color.rgb = RGBColor(150, 150, 150)
        placeholder_shape.line.width = Pt(1)


class StabilityReportPipelineV2:
    """稳定性报告生成管道 v2.0 - 两阶段架构"""
    
    def __init__(self, 
                 output_dir: str = "reports",
                 assets_dir: Optional[str] = None,
                 config_path: str = 'catalog_config.yaml',
                 overwrite_mode: bool = False):
        """
        初始化报告生成管道
        
        Args:
            output_dir: 输出目录
            assets_dir: 素材目录（如果不指定，会在output_dir下创建assets子目录）
            config_path: catalog配置文件路径
            overwrite_mode: 是否覆盖已存在的素材文件。True=覆盖，False=跳过已存在的文件（默认）
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # 素材目录
        if assets_dir:
            self.assets_dir = Path(assets_dir)
        else:
            self.assets_dir = self.output_dir / "assets"
        
        # 初始化组件
        self.asset_manager = AssetManager(self.assets_dir)
        self.asset_generator = AssetGenerator(self.asset_manager, config_path, overwrite_mode)
        self.ppt_composer = PPTComposer(self.asset_manager)
        
        logger.info(f"Pipeline v2.0 initialized. Output: {self.output_dir}, Assets: {self.assets_dir}")
    
    def generate_chip_report(self, chip_id: str, skip_asset_generation: bool = False) -> str:
        """
        为指定chip生成完整报告
        
        Args:
            chip_id: chip ID
            skip_asset_generation: 是否跳过素材生成（用于测试PPT布局）
            
        Returns:
            str: PPT文件路径
        """
        logger.info(f"Starting report generation for chip {chip_id}")
        
        # 阶段1：生成素材（如果需要）
        if not skip_asset_generation:
            logger.info(f"Phase 1: Generating assets for chip {chip_id}")
            asset_results = self.asset_generator.generate_chip_assets(chip_id)
            
            if not asset_results['success']:
                logger.error(f"Asset generation failed for chip {chip_id}")
                return None
            
            logger.info(f"Asset generation completed: {asset_results['generated_assets']}")
        else:
            logger.info("Skipping asset generation (using existing assets)")
        
        # 阶段2：创建PPT
        logger.info(f"Phase 2: Creating PPT presentation for chip {chip_id}")
        
        output_filename = f"stability_report_{chip_id.replace('#', 'chip')}.pptx"
        output_path = self.output_dir / output_filename
        
        final_ppt = self.ppt_composer.create_presentation(chip_id, output_path)
        
        if final_ppt:
            logger.info(f"Report generated successfully: {final_ppt}")
            return final_ppt
        else:
            logger.error("Failed to create PPT presentation")
            return None
    
    def generate_all_chip_reports(self, skip_asset_generation: bool = False) -> List[str]:
        """
        为所有chip生成报告
        
        Args:
            skip_asset_generation: 是否跳过素材生成
            
        Returns:
            List[str]: 生成的PPT文件路径列表
        """
        logger.info("Starting batch report generation for all chips")
        
        # 发现所有可用的chip
        try:
            all_experiments = self.asset_generator.manager.search()
            chip_ids = list(set(exp.chip_id for exp in all_experiments))
            chip_ids.sort()
            
            logger.info(f"Found {len(chip_ids)} total chips: {chip_ids[:5]}...")
            
            generated_reports = []
            for i, chip_id in enumerate(chip_ids, 1):
                logger.info(f"Processing chip {chip_id} ({i}/{len(chip_ids)})")
                
                try:
                    report_path = self.generate_chip_report(chip_id, skip_asset_generation)
                    if report_path:
                        generated_reports.append(report_path)
                        logger.info(f"Successfully generated report: {report_path}")
                    else:
                        logger.warning(f"Failed to generate report for chip {chip_id}")
                        
                except Exception as e:
                    logger.error(f"Error processing chip {chip_id}: {e}")
            
            logger.info(f"Batch generation completed. Generated {len(generated_reports)} reports out of {len(chip_ids)} chips")
            return generated_reports
            
        except Exception as e:
            logger.error(f"Failed to generate batch reports: {e}")
            return []


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description="OECT器件稳定性分析报告生成器 v2.0")
    
    parser.add_argument('--chip', type=str, help='指定chip ID，如 "#20250804008"')
    parser.add_argument('--all-chips', action='store_true', help='为所有chip生成报告')
    parser.add_argument('--output', type=str, help='输出PPT文件路径（单个chip）')
    parser.add_argument('--output-dir', type=str, default='reports', help='输出目录（多个chip）')
    parser.add_argument('--assets-dir', type=str, help='素材目录（可选）')
    parser.add_argument('--config', type=str, default='catalog_config.yaml', help='Catalog配置文件路径')
    parser.add_argument('--skip-assets', action='store_true', help='跳过素材生成（仅用于测试PPT布局）')
    parser.add_argument('--overwrite', action='store_true', help='覆盖已存在的素材文件（默认为跳过）')
    parser.add_argument('--verbose', action='store_true', help='详细日志输出')
    
    args = parser.parse_args()
    
    # 配置日志级别
    if args.verbose:
        logging.basicConfig(level=logging.DEBUG)
    else:
        logging.basicConfig(level=logging.INFO)
    
    try:
        # 初始化管道
        pipeline = StabilityReportPipelineV2(
            output_dir=args.output_dir,
            assets_dir=args.assets_dir,
            config_path=args.config,
            overwrite_mode=args.overwrite
        )
        
        if args.all_chips:
            # 批量生成所有chip的报告
            logger.info("Generating reports for all chips")
            reports = pipeline.generate_all_chip_reports(skip_asset_generation=args.skip_assets)
            
            if reports:
                logger.info(f"Successfully generated {len(reports)} reports")
                for report in reports:
                    print(f"Generated: {report}")
            else:
                logger.error("No reports were generated")
                
        elif args.chip:
            # 生成指定chip的报告
            logger.info(f"Generating report for chip {args.chip}")
            
            if args.output:
                # 指定了输出路径，直接使用PPT编排器
                asset_manager = AssetManager(pipeline.assets_dir)
                ppt_composer = PPTComposer(asset_manager)
                report_path = ppt_composer.create_presentation(args.chip, args.output)
            else:
                report_path = pipeline.generate_chip_report(args.chip, skip_asset_generation=args.skip_assets)
            
            if report_path:
                logger.info(f"Report generated successfully: {report_path}")
                print(f"Generated: {report_path}")
            else:
                logger.error("Failed to generate report")
        else:
            parser.print_help()
            
    except Exception as e:
        logger.error(f"Pipeline execution failed: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()