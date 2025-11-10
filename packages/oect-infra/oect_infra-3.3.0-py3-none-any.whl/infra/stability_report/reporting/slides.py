from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

try:
    import yaml
except ImportError as exc:  # pragma: no cover
    raise RuntimeError("PyYAML is required for slide configuration") from exc

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .assets import AssetManifest, ManifestStore

logger = logging.getLogger(__name__)


@dataclass
class SlideDefinition:
    id: str
    builder: str
    params: Dict[str, Any] = field(default_factory=dict)
    scope: Dict[str, Any] = field(default_factory=dict)
    enabled: bool = True

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> "SlideDefinition":
        return cls(
            id=data["id"],
            builder=data["builder"],
            params=data.get("params", {}),
            scope=data.get("scope", {}),
            enabled=data.get("enabled", True),
        )


class SlideRegistry:
    def __init__(self) -> None:
        self._builders: Dict[str, type[SlideBuilder]] = {}

    def register(self, name: str, builder_cls: type["SlideBuilder"]) -> None:
        self._builders[name] = builder_cls

    def get(self, name: str) -> type["SlideBuilder"]:
        if name not in self._builders:
            raise KeyError(f"No slide builder registered for '{name}'")
        return self._builders[name]


@dataclass
class SlideContext:
    chip_id: str
    manifest: AssetManifest
    assets_dir: Path
    slides_dir: Path
    params: Dict[str, Any]


class SlideBuilder:
    def __init__(self, params: Dict[str, Any]):
        self.params = params

    def render(self, presentation: Presentation, context: SlideContext) -> None:
        raise NotImplementedError

    @staticmethod
    def _blank_slide(prs: Presentation):
        return prs.slides.add_slide(prs.slide_layouts[6])


class ImageGridSlide(SlideBuilder):
    def render(self, presentation: Presentation, context: SlideContext) -> None:
        slide = self._blank_slide(presentation)
        title_text = self.params.get("title", "")
        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
            frame = title_box.text_frame
            frame.text = title_text.format(chip_id=context.chip_id)
            frame.paragraphs[0].font.size = Pt(20)
            frame.paragraphs[0].font.bold = True
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        asset_id = self.params["asset_id"]
        device_ids = self.params.get("device_ids") or [rec.device_id for rec in context.manifest.records if rec.asset_id == asset_id]
        device_ids = [d for d in device_ids if d]
        columns = self.params.get("columns", 3)
        padding = self.params.get("padding", 0.1)
        width = Inches(12.33 - padding * 2)
        height = Inches(6.0)
        start_x = Inches(padding)
        start_y = Inches(1.2)
        cell_w = width / columns
        rows = max(1, ((len(device_ids) - 1) // columns) + 1)
        cell_h = height / rows

        for idx, device_id in enumerate(device_ids):
            row = idx // columns
            col = idx % columns
            x = start_x + col * cell_w
            y = start_y + row * cell_h
            records = context.manifest.find(asset_id, device_id=device_id)
            if records:
                image_path = Path(records[0].path)
                if image_path.exists():
                    slide.shapes.add_picture(
                        str(image_path),
                        x + Inches(0.05),
                        y + Inches(0.05),
                        width=cell_w - Inches(0.1),
                        height=cell_h - Inches(0.3),
                    )
                    caption = slide.shapes.add_textbox(x, y + cell_h - Inches(0.25), cell_w, Inches(0.2))
                    frame = caption.text_frame
                    frame.text = self.params.get("caption_template", "Device {device_id}").format(device_id=device_id)
                    frame.paragraphs[0].font.size = Pt(11)
                    frame.paragraphs[0].font.bold = True
                    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    continue
            placeholder = slide.shapes.add_textbox(x, y, cell_w, cell_h)
            frame = placeholder.text_frame
            frame.text = self.params.get("fallback_text", "No Data")
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            frame.paragraphs[0].font.size = Pt(12)


class FeatureMatrixSlide(SlideBuilder):
    def render(self, presentation: Presentation, context: SlideContext) -> None:
        slide = self._blank_slide(presentation)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        frame = title_box.text_frame
        frame.text = self.params.get("title", "Feature Analysis").format(chip_id=context.chip_id)
        frame.paragraphs[0].font.size = Pt(20)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        asset_id = self.params["asset_id"]
        variants = self.params.get("variants", [])
        columns = self.params.get("columns", 3)
        rows = max(1, ((len(variants) - 1) // columns) + 1)
        padding = self.params.get("padding", 0.1)
        start_x = Inches(padding)
        start_y = Inches(1.2)
        width = Inches(12.33 - padding * 2)
        height = Inches(6.0)
        cell_w = width / columns
        cell_h = height / rows

        for idx, variant_cfg in enumerate(variants):
            row = idx // columns
            col = idx % columns
            x = start_x + col * cell_w
            y = start_y + row * cell_h
            variant = variant_cfg.get("variant")
            records = context.manifest.find(asset_id, variant=variant)
            if records:
                image_path = Path(records[0].path)
                if image_path.exists():
                    slide.shapes.add_picture(
                        str(image_path),
                        x + Inches(0.05),
                        y + Inches(0.05),
                        width=cell_w - Inches(0.1),
                        height=cell_h - Inches(0.35),
                    )
                    caption = slide.shapes.add_textbox(x, y + cell_h - Inches(0.25), cell_w, Inches(0.2))
                    caption_frame = caption.text_frame
                    caption_frame.text = variant_cfg.get("title", variant or "Variant")
                    caption_frame.paragraphs[0].font.size = Pt(10)
                    caption_frame.paragraphs[0].font.bold = True
                    caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    continue
            placeholder = slide.shapes.add_textbox(x, y, cell_w, cell_h)
            placeholder.text_frame.text = self.params.get("fallback_text", "No Data")
            placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            placeholder.text_frame.paragraphs[0].font.size = Pt(12)


class VideoGridSlide(SlideBuilder):
    def render(self, presentation: Presentation, context: SlideContext) -> None:
        slide = self._blank_slide(presentation)
        title_text = self.params.get("title", "")
        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
            frame = title_box.text_frame
            frame.text = title_text.format(chip_id=context.chip_id)
            frame.paragraphs[0].font.size = Pt(20)
            frame.paragraphs[0].font.bold = True
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        asset_id = self.params["asset_id"]
        device_ids = self.params.get("device_ids") or [rec.device_id for rec in context.manifest.records if rec.asset_id == asset_id]
        device_ids = [d for d in device_ids if d]
        columns = self.params.get("columns", 3)
        padding = self.params.get("padding", 0.1)
        start_x = Inches(padding)
        start_y = Inches(1.2)
        width = Inches(12.33 - padding * 2)
        height = Inches(6.0)
        rows = max(1, ((len(device_ids) - 1) // columns) + 1)
        cell_w = width / columns
        cell_h = height / rows

        for idx, device_id in enumerate(device_ids):
            row = idx // columns
            col = idx % columns
            x = start_x + col * cell_w
            y = start_y + row * cell_h
            records = context.manifest.find(asset_id, device_id=device_id)
            if records:
                video_path = Path(records[0].path)
                if video_path.exists():
                    slide.shapes.add_movie(
                        str(video_path),
                        x + Inches(0.05),
                        y + Inches(0.05),
                        width=cell_w - Inches(0.1),
                        height=cell_h - Inches(0.35),
                        poster_frame_image=None,
                        mime_type="video/mp4",
                    )
                    caption = slide.shapes.add_textbox(x, y + cell_h - Inches(0.25), cell_w, Inches(0.2))
                    frame = caption.text_frame
                    frame.text = self.params.get("caption_template", "Device {device_id}").format(device_id=device_id)
                    frame.paragraphs[0].font.size = Pt(10)
                    frame.paragraphs[0].font.bold = True
                    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    continue
            placeholder = slide.shapes.add_textbox(x, y, cell_w, cell_h)
            placeholder.text_frame.text = self.params.get("fallback_text", "Missing Video")
            placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            placeholder.text_frame.paragraphs[0].font.size = Pt(12)


class SlideEngine:
    def __init__(
        self,
        assets_dir: Path,
        slides_dir: Path,
        config_path: Path,
        registry: Optional[SlideRegistry] = None,
    ) -> None:
        self.assets_dir = Path(assets_dir)
        self.slides_dir = Path(slides_dir)
        self.slides_dir.mkdir(parents=True, exist_ok=True)
        self.config_path = Path(config_path)
        self.registry = registry or default_slide_registry()
        self.manifest_store = ManifestStore(self.assets_dir / "manifests")
        self.slide_definitions = self._load_config()

    def _load_config(self) -> List[SlideDefinition]:
        config = yaml.safe_load(self.config_path.read_text())
        if not config or "slides" not in config:
            raise ValueError(f"Slide configuration '{self.config_path}' missing 'slides' list")
        return [SlideDefinition.from_dict(item) for item in config["slides"] if item.get("enabled", True)]

    def build_all(self) -> List[Path]:
        return self.build()

    def build(
        self,
        *,
        slide_ids: Optional[Iterable[str]] = None,
        chips: Optional[Iterable[str]] = None,
    ) -> List[Path]:
        outputs: List[Path] = []
        available_chips = self.manifest_store.list_chips()
        if chips is not None:
            target_chips = [c for c in chips if c in available_chips]
        else:
            target_chips = available_chips
        slide_map = {definition.id: definition for definition in self.slide_definitions}
        if slide_ids is None:
            selected = list(self.slide_definitions)
        else:
            selected = [slide_map[sid] for sid in slide_ids if sid in slide_map]
        for slide_def in selected:
            chips_for_slide = self._chips_for_definition(slide_def, target_chips)
            for chip_id in chips_for_slide:
                manifest = self.manifest_store.load(chip_id)
                if not manifest:
                    logger.warning("Manifest missing for chip %s", chip_id)
                    continue
                path = self._build_slide(slide_def, chip_id, manifest)
                if path:
                    outputs.append(path)
        return outputs

    def _chips_for_definition(self, slide_def: SlideDefinition, available: Iterable[str]) -> List[str]:
        scope = slide_def.scope or {}
        foreach = scope.get("foreach", "chip")
        if foreach == "none":
            chip = scope.get("chip")
            if not chip:
                raise ValueError(f"Slide '{slide_def.id}' scope requires 'chip'")
            return [chip]
        chips = list(available)
        allowed = scope.get("chips")
        if allowed:
            normalized = []
            for item in allowed:
                if item == "*":
                    normalized.extend(chips)
                elif item in chips:
                    normalized.append(item)
            if normalized:
                chips = sorted(set(normalized))
        return chips

    def _build_slide(self, slide_def: SlideDefinition, chip_id: str, manifest: AssetManifest) -> Optional[Path]:
        builder_cls = self.registry.get(slide_def.builder)
        builder = builder_cls(slide_def.params)
        presentation = Presentation()
        presentation.slide_width = Inches(13.33)
        presentation.slide_height = Inches(7.5)
        context = SlideContext(
            chip_id=chip_id,
            manifest=manifest,
            assets_dir=self.assets_dir,
            slides_dir=self.slides_dir,
            params=slide_def.params,
        )
        builder.render(presentation, context)
        slide_dir = self.slides_dir / slide_def.id
        slide_dir.mkdir(parents=True, exist_ok=True)
        output_path = slide_dir / f"{chip_id.replace('#', 'chip').replace('/', '_')}.pptx"
        presentation.save(str(output_path))
        self._write_metadata(slide_dir / f"{output_path.stem}.json", slide_def, chip_id, manifest)
        logger.info("Built slide %s for chip %s -> %s", slide_def.id, chip_id, output_path)
        return output_path

    def _write_metadata(self, path: Path, slide_def: SlideDefinition, chip_id: str, manifest: AssetManifest) -> None:
        data = {
            "slide_id": slide_def.id,
            "chip_id": chip_id,
            "builder": slide_def.builder,
            "params": slide_def.params,
            "generated_at": datetime.utcnow().isoformat(),
            "assets": [rec.to_dict() for rec in manifest.records if rec.asset_id == slide_def.params.get("asset_id")],
        }
        path.write_text(json.dumps(data, indent=2, ensure_ascii=True))


def default_slide_registry() -> SlideRegistry:
    registry = SlideRegistry()
    registry.register("image_grid", ImageGridSlide)
    registry.register("feature_matrix", FeatureMatrixSlide)
    registry.register("video_grid", VideoGridSlide)
    return registry
