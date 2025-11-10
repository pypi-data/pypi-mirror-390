from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    import yaml
except ImportError as exc:  # pragma: no cover
    raise RuntimeError("PyYAML is required for report configuration") from exc

from pptx import Presentation
from pptx.util import Inches

from .assets import ManifestStore
from .slides import (
    SlideContext,
    SlideDefinition,
    SlideEngine,
    SlideRegistry,
    default_slide_registry,
)

logger = logging.getLogger(__name__)


@dataclass
class ReportSlide:
    slide_id: str
    chips: List[str]
    params: Dict[str, Any] = field(default_factory=dict)

    @classmethod
    def from_dict(cls, data: Dict[str, Any], *, available_chips: List[str]) -> "ReportSlide":
        chips_spec = data.get("chips")
        if chips_spec is None:
            raise ValueError(f"Report slide '{data.get('slide_id')}' requires 'chips' list")
        chips = _normalize_chips(chips_spec, available_chips)
        return cls(slide_id=data["slide_id"], chips=chips, params=data.get("params", {}))


def _normalize_chips(spec: Any, available: List[str]) -> List[str]:
    if spec == "*":
        return list(available)
    if isinstance(spec, str):
        return [spec]
    if isinstance(spec, list):
        result: List[str] = []
        for item in spec:
            if item == "*":
                result.extend(available)
            elif item in available:
                result.append(item)
            else:
                result.append(item)
        return list(dict.fromkeys(result))
    raise ValueError(f"Invalid chips specification: {spec!r}")


@dataclass
class ReportDefinition:
    id: str
    title: Optional[str] = None
    theme: Optional[str] = None
    slides: List[ReportSlide] = field(default_factory=list)

    @classmethod
    def from_dict(cls, data: Dict[str, Any], *, available_chips: List[str]) -> "ReportDefinition":
        slides = [ReportSlide.from_dict(item, available_chips=available_chips) for item in data.get("slides", [])]
        return cls(
            id=data["id"],
            title=data.get("title"),
            theme=data.get("theme"),
            slides=slides,
        )


class ReportComposer:
    def __init__(
        self,
        assets_dir: Path,
        slides_config_path: Path,
        reports_config_path: Path,
        output_dir: Path,
        registry: Optional[SlideRegistry] = None,
    ) -> None:
        self.assets_dir = Path(assets_dir)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.manifest_store = ManifestStore(self.assets_dir / "manifests")
        registry = registry or default_slide_registry()
        self.registry = registry
        self.slide_engine = SlideEngine(
            assets_dir=self.assets_dir,
            slides_dir=self.assets_dir / "_tmp_slides",
            config_path=slides_config_path,
            registry=registry,
        )
        self.report_definitions = self._load_reports(reports_config_path)
        self.slides_lookup = {definition.id: definition for definition in self.slide_engine.slide_definitions}

    def _load_reports(self, config_path: Path) -> Dict[str, ReportDefinition]:
        config = yaml.safe_load(Path(config_path).read_text())
        if not config or "reports" not in config:
            raise ValueError(f"Report configuration '{config_path}' missing 'reports' list")
        chips = self.manifest_store.list_chips()
        definitions = [ReportDefinition.from_dict(item, available_chips=chips) for item in config["reports"]]
        return {definition.id: definition for definition in definitions}

    def compose(self, report_id: str, *, output_path: Optional[Path] = None) -> Path:
        if report_id not in self.report_definitions:
            raise KeyError(f"Unknown report id '{report_id}'")
        definition = self.report_definitions[report_id]
        presentation = self._create_presentation(definition)
        for slide_ref in definition.slides:
            slide_def = self.slides_lookup.get(slide_ref.slide_id)
            if not slide_def:
                logger.warning("Slide definition '%s' not found", slide_ref.slide_id)
                continue
            builder_cls = self.registry.get(slide_def.builder)
            base_params = dict(slide_def.params)
            base_params.update(slide_ref.params)
            builder = builder_cls(base_params)
            for chip_id in slide_ref.chips:
                manifest = self.manifest_store.load(chip_id)
                if not manifest:
                    logger.warning("Manifest missing for chip %s", chip_id)
                    continue
                context = SlideContext(
                    chip_id=chip_id,
                    manifest=manifest,
                    assets_dir=self.assets_dir,
                    slides_dir=self.output_dir,
                    params=base_params,
                )
                builder.render(presentation, context)
        target = output_path or (self.output_dir / f"{report_id}.pptx")
        presentation.save(str(target))
        logger.info("Report '%s' generated -> %s", report_id, target)
        return target

    def list_reports(self) -> List[str]:
        return list(self.report_definitions.keys())

    @staticmethod
    def _create_presentation(definition: ReportDefinition) -> Presentation:
        if definition.theme:
            template = Path(definition.theme)
            if template.exists():
                prs = Presentation(str(template))
            else:
                logger.warning("Theme '%s' not found, using default", definition.theme)
                prs = Presentation()
        else:
            prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        if definition.title:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = definition.title
        return prs
