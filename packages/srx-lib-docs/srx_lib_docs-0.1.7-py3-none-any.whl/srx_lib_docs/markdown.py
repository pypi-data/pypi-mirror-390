from __future__ import annotations

import os
import aiofiles
import httpx
import tempfile
import logging
from markitdown import MarkItDown
from typing import Dict, Any, Optional

logger = logging.getLogger(__name__)


class DocumentMarkdownConverter:
    """Convert various document formats to markdown (async helpers)."""

    def __init__(self):
        self.supported_types = {
            "application/pdf": "pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
            "application/msword": "doc",
            "application/vnd.ms-powerpoint": "ppt",
            # Audio formats for transcription
            "audio/mpeg": "mp3",
            "audio/mp3": "mp3",
            "audio/x-m4a": "m4a",
            "audio/m4a": "m4a",
            "audio/wav": "wav",
            "audio/x-wav": "wav",
        }

    async def download_file(self, url: str, temp_dir: str = "/tmp") -> str:
        os.makedirs(temp_dir, exist_ok=True)
        tmp = tempfile.NamedTemporaryFile(delete=False, dir=temp_dir, suffix=".tmp")
        temp_path = tmp.name
        tmp.close()

        logger.info("Downloading file from %s", url)
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                async with client.stream("GET", url) as response:
                    response.raise_for_status()
                    async with aiofiles.open(temp_path, "wb") as f:
                        async for chunk in response.aiter_bytes():
                            await f.write(chunk)
            logger.info("File downloaded to %s", temp_path)
            return temp_path
        except Exception:
            if os.path.exists(temp_path):
                try:
                    os.unlink(temp_path)
                except Exception:
                    pass
            raise

    def get_file_type_from_mimetype(self, mimetype: str) -> str:
        if mimetype in self.supported_types:
            return self.supported_types[mimetype]
        raise ValueError(f"Unsupported mimetype: {mimetype}")

    async def convert_pdf_to_markdown(self, file_path: str) -> str:
        import PyPDF2

        parts: list[str] = ["# PDF Document\n\n"]
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(f"## Page {i}\n\n")
                    parts.append(text)
                    parts.append("\n\n")
        return "".join(parts)

    async def convert_docx_to_markdown(self, file_path: str) -> str:
        from docx import Document

        md: list[str] = []
        doc = Document(file_path)
        for p in doc.paragraphs:
            text = (p.text or "").strip()
            if not text:
                continue
            if p.style and p.style.name and p.style.name.startswith("Heading"):
                level = p.style.name[-1] if p.style.name[-1].isdigit() else "1"
                md.append(f"{'#' * int(level)} {text}\n\n")
            else:
                md.append(f"{text}\n\n")
        for table in doc.tables:
            md.append("\n")
            for i, row in enumerate(table.rows):
                cells = [c.text.strip() for c in row.cells]
                if i == 0:
                    md.append("| " + " | ".join(cells) + " |\n")
                    md.append("| " + " | ".join(["---"] * len(cells)) + " |\n")
                else:
                    md.append("| " + " | ".join(cells) + " |\n")
            md.append("\n")
        return "".join(md)

    async def convert_pptx_to_markdown(self, file_path: str) -> str:
        from pptx import Presentation

        md: list[str] = ["# PowerPoint Presentation\n\n"]
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            md.append(f"## Slide {i}\n\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    md.append(f"{shape.text}\n\n")
        return "".join(md)

    async def convert_xlsx_to_markdown(self, file_path: str) -> str:
        from openpyxl import load_workbook
        import warnings

        md: list[str] = ["# Excel Spreadsheet\n\n"]
        with warnings.catch_warnings():
            warnings.filterwarnings(
                "ignore",
                category=UserWarning,
                message=".*Data Validation extension is not supported.*",
            )
            wb = load_workbook(file_path, data_only=True, read_only=True, keep_links=False)

        for sheet_name in wb.sheetnames:
            md.append(f"## Sheet: {sheet_name}\n\n")
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                if not row:
                    continue
                cells = ["" if v is None else str(v) for v in row]
                if any((c or "").strip() for c in cells):
                    md.append("| " + " | ".join(cells) + " |\n")
            md.append("\n")
        return "".join(md)

    async def convert_to_markdown(self, file_path: str, file_type: str) -> str:
        """Convert using MarkItDown as the primary path; fall back to simple converters.

        For audio files (mp3, m4a, wav), MarkItDown uses speech_recognition for transcription.
        Note: Audio transcription requires the 'audio-transcription' extras to be installed.
        """
        try:
            md = MarkItDown(enable_plugins=False)
            res = md.convert(file_path)
            text = getattr(res, "text_content", None)
            if isinstance(text, str) and text.strip():
                return text
        except Exception as e:
            logger.warning("MarkItDown convert failed (%s); falling back", e)

        # Fallback per-format minimal converters (audio files have no fallback)
        if file_type == "pdf":
            return await self.convert_pdf_to_markdown(file_path)
        if file_type == "docx":
            return await self.convert_docx_to_markdown(file_path)
        if file_type == "pptx":
            return await self.convert_pptx_to_markdown(file_path)
        if file_type == "xlsx":
            return await self.convert_xlsx_to_markdown(file_path)
        if file_type in ("mp3", "m4a", "wav"):
            raise ValueError(
                f"Audio file transcription failed. Ensure 'markitdown[audio-transcription]' "
                f"is installed and ffmpeg is available on the system."
            )
        raise ValueError(f"Unsupported file type: {file_type}")

    async def process_document(
        self, url: str, metadata: Optional[Dict[str, Any]] = None, mimetype: Optional[str] = None
    ) -> Dict[str, Any]:
        if not mimetype:
            raise ValueError("Mimetype is required")
        temp_file_path = None
        try:
            temp_file_path = await self.download_file(url)
            file_type = self.get_file_type_from_mimetype(mimetype)
            markdown_content = await self.convert_to_markdown(temp_file_path, file_type)
            file_size = os.path.getsize(temp_file_path)
            logger.info("Document processed: %s, %s bytes", file_type, file_size)
            return {
                "success": True,
                "file_type": file_type,
                "file_size": file_size,
                "markdown_content": markdown_content,
                "metadata": metadata or {},
            }
        except Exception as e:
            logger.error("Document processing failed: %s", e)
            return {
                "success": False,
                "error": str(e),
                "file_type": "unknown",
                "metadata": metadata or {},
            }
        finally:
            if temp_file_path and os.path.exists(temp_file_path):
                try:
                    os.unlink(temp_file_path)
                    logger.info("Cleaned up temp file: %s", temp_file_path)
                except Exception as e:
                    logger.warning("Failed to clean up temp file %s: %s", temp_file_path, e)


__all__ = ["DocumentMarkdownConverter"]
