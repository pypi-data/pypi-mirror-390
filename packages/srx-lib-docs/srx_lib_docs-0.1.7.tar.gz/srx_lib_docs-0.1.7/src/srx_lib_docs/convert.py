from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Optional, Union

from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from PyPDF2 import PdfReader


PathOrBytes = Union[str, Path, bytes, bytearray, BytesIO]


def _to_bytes_io(inp: PathOrBytes) -> BytesIO:
    if isinstance(inp, (bytes, bytearray)):
        return BytesIO(inp)
    if isinstance(inp, BytesIO):
        return inp
    p = Path(str(inp))
    return BytesIO(p.read_bytes())


def extract_text(path_or_bytes: PathOrBytes, mime_type: Optional[str] = None) -> str:
    """Extract plain text from PDF, DOCX, PPTX, XLSX.

    - If `mime_type` is not provided, infer from file suffix when a path is given.
    - Light-weight logic intended for search/index or LLM context building.
    """
    guessed = None
    if mime_type:
        guessed = mime_type.lower()
    else:
        if not isinstance(path_or_bytes, (bytes, bytearray, BytesIO)):
            suf = Path(str(path_or_bytes)).suffix.lower()
            guessed = {
                ".pdf": "application/pdf",
                ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            }.get(suf)

    if guessed in ("application/pdf",) or (not guessed and str(path_or_bytes).lower().endswith(".pdf")):
        bio = _to_bytes_io(path_or_bytes)
        reader = PdfReader(bio)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if guessed and "wordprocessingml" in guessed or (
        not guessed and str(path_or_bytes).lower().endswith(".docx")
    ):
        bio = _to_bytes_io(path_or_bytes)
        doc = DocxDocument(bio)
        return "\n".join(p.text or "" for p in doc.paragraphs)

    if guessed and "presentationml" in guessed or (
        not guessed and str(path_or_bytes).lower().endswith(".pptx")
    ):
        bio = _to_bytes_io(path_or_bytes)
        prs = Presentation(bio)
        texts = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    texts.append(shp.text)
        return "\n".join(texts)

    if guessed and "spreadsheetml" in guessed or (
        not guessed and str(path_or_bytes).lower().endswith(".xlsx")
    ):
        bio = _to_bytes_io(path_or_bytes)
        wb = load_workbook(bio, data_only=True, read_only=True)
        lines = []
        for ws in wb.worksheets:
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i > 20:
                    break
                if not row:
                    lines.append("")
                else:
                    lines.append(
                        ", ".join(str(c) if c is not None else "" for c in row)
                    )
        return "\n".join(lines)

    # Fallback: if it's bytes and no type, return empty string
    return ""
