"""
File Reader Module
==================

This module provides a unified interface for reading various file formats
into plain text. It supports documents, data files, presentations, code files,
and more.

Supported Formats:
    - Documents: PDF, DOCX, HTML, TXT, Markdown, RTF
    - Data: CSV, TSV, Excel (XLS/XLSX), JSON, XML, YAML
    - Presentations: PPTX
    - Code: Python, JavaScript, TypeScript, Java, C++, Go, Rust, PHP, etc.
    - Config: INI, TOML, ENV, properties
    - Logs: LOG, SRT, VTT

The FileReader class automatically detects file types by extension and
applies the appropriate reader. Unsupported formats fall back to plain
text reading.

Example:
    >>> from neurosurfer.rag.filereader import FileReader
    >>> 
    >>> reader = FileReader()
    >>> 
    >>> # Read PDF
    >>> text = reader.read("document.pdf")
    >>> 
    >>> # Read Excel
    >>> text = reader.read("data.xlsx")
    >>> 
    >>> # Read Python code
    >>> text = reader.read("script.py")
"""
import os
from pathlib import Path
import fitz  # for PDFs
import docx  # for DOCX
import pandas as pd  # for Excel and CSV
from bs4 import BeautifulSoup  # for HTML

# Optional dependencies (handle if not installed)
try:
    import yaml
except ImportError:
    yaml = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import xml.etree.ElementTree as ET
except ImportError:
    ET = None
    

class FileReader:
    """
    Unified file reader for multiple formats.
    
    This class provides a single interface for reading various file types
    into plain text. It automatically detects the file type by extension
    and applies the appropriate reader method.
    
    Attributes:
        supported_types (dict): Mapping of file extensions to reader functions
    
    Example:
        >>> reader = FileReader()
        >>> 
        >>> # Read different file types
        >>> pdf_text = reader.read("report.pdf")
        >>> excel_text = reader.read("data.xlsx")
        >>> code_text = reader.read("script.py")
        >>> 
        >>> # Check supported types
        >>> print(reader.supported_types.keys())
    """
    def __init__(self):
        self.supported_types = {
            # General
            ".pdf": self._read_pdf,
            ".html": self._read_html,
            ".htm": self._read_html,
            ".docx": self._read_docx,
            
            # Data
            ".csv": self._read_csv,
            ".tsv": self._read_csv,
            ".xls": self._read_excel,
            ".xlsx": self._read_excel,
            ".xml": self._read_xml,
            ".yaml": self._read_yaml,
            ".yml": self._read_yaml,

            # Presentations
            ".pptx": self._read_pptx,
        }
        # Add Code files, Config, Misc etc which can be read as normal txt to the supported files
        read_as_txt_lits = [
            ".txt", ".md", ".rtf", ".doc", ".odt", ".json", ".ppt", ".py", ".ipynb", ".java", ".js", ".ts", ".jsx", ".tsx", ".cpp", ".c", ".h", 
            ".cs", ".go", ".rb", ".rs", ".php", ".swift", ".kt", ".sh", ".bat", ".ps1", ".scala", ".lua", ".r", ".env", ".ini", ".toml", ".cfg", 
            ".conf", ".properties", ".log", ".tex", ".srt", ".vtt",
        ]
        for ext in read_as_txt_lits:
            self.supported_types[ext] = self._read_txt

    def read(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        reader_fn = self.supported_types.get(ext, self._read_txt)
        return reader_fn(file_path)

    def _read_pdf(self, path: str) -> str:
        text = ""
        try:
            with fitz.open(path) as doc:
                for page in doc:
                    text += page.get_text()
        except Exception as e:
            text = f"Error reading PDF: {e}"
        return text

    def _read_txt(self, path: str) -> str:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            return f"Error reading TXT: {e}"

    def _read_html(self, path: str) -> str:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f, "html.parser")
                return soup.get_text()
        except Exception as e:
            return f"Error reading HTML: {e}"

    def _read_docx(self, path: str) -> str:
        try:
            doc = docx.Document(path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            return f"Error reading DOCX: {e}"

    def _read_excel(self, path: str) -> str:
        try:
            text = ""
            xls = pd.read_excel(path, sheet_name=None)
            for sheet_name, sheet_data in xls.items():
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += sheet_data.astype(str).to_string(index=False)
            return text
        except Exception as e:
            return f"Error reading Excel: {e}"

    def _read_csv(self, path: str) -> str:
        try:
            df = pd.read_csv(path)
            return df.astype(str).to_string(index=False)
        except Exception as e:
            return f"Error reading CSV/TSV: {e}"

    def _read_yaml(self, path: str) -> str:
        if not yaml:
            return "PyYAML not installed"
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = yaml.safe_load(f)
                return str(data)
        except Exception as e:
            return f"Error reading YAML: {e}"

    def _read_xml(self, path: str) -> str:
        if not ET:
            return "XML parser not available"
        try:
            tree = ET.parse(path)
            root = tree.getroot()
            return ET.tostring(root, encoding="unicode")
        except Exception as e:
            return f"Error reading XML: {e}"

    def _read_pptx(self, path: str) -> str:
        if not Presentation:
            return "python-pptx not installed"
        try:
            prs = Presentation(path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            return f"Error reading PPTX: {e}"
